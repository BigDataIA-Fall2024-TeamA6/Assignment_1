import os
import streamlit as st
from openai import OpenAI
import cv2
import pandas as pd
import boto3
from dotenv import load_dotenv
from db import DBConnection
from pydub import AudioSegment
import speech_recognition as sr
from io import BytesIO
 
import PyPDF2
from pptx import Presentation
from docx import Document
import zipfile
from PIL import Image
import easyocr
import pytesseract
 
import tiktoken  # OpenAI's token counter library
 
load_dotenv()
# Securely get environment variables
api_key = os.getenv("OPENAI_API_KEY")  # OpenAI API key
aws_access_key_id = os.getenv("AWS_ACCESS_KEY_ID")  # AWS access key
aws_secret_access_key = os.getenv("AWS_SECRET_ACCESS_KEY")  # AWS secret key
aws_region = os.getenv("AWS_REGION", "us-east-1")  # AWS region, default to us-east-1
 
# Initialize OpenAI client with the API key
client = OpenAI(api_key=api_key)
 
# Initialize Boto3 client for S3
s3_client = boto3.client(
    's3',
    aws_access_key_id=aws_access_key_id,
    aws_secret_access_key=aws_secret_access_key,
    region_name=aws_region
)
 
# Code to extract text from audio
 
 
# code to extract text from PDF
 
def extract_text_from_pptx(content):
    presentation = Presentation(BytesIO(content))
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text
 
# code to extract text from docx
 
def extract_text_from_image(local_path):
    try:
        reader = easyocr.Reader(['en'],gpu=True)
        img = cv2.imread(local_path)
        result = reader.readtext(img)
        text=""
        for detection in result:
            text += " " + detection[1]
        return text
   
    except Exception as e:
        st.error(f"Error processing image for OCR: {e}")
        return None
   
def analyze_image(local_path):
    try:
        text = extract_text_from_image(local_path)
        if text:
            description += f"Here is th text from the image: '{text}'"
        else:
            description = "No readable text was detected in the image."
        return description
    except Exception as e:
        st.error(f"Error processing image: {e}")
        return None
 
 
def extract_text_from_zip(content):
    with zipfile.ZipFile(BytesIO(content), 'r') as z:
        text = ""
        for file_name in z.namelist():
            with z.open(file_name) as f:
                text += f.read().decode('utf-8', errors='ignore')
    return text
 
 
def count_tokens(text, model="gpt-4"):    
    # Count the number of tokens in the text using the OpenAI tokenization system.
    encoding = tiktoken.encoding_for_model(model)
    return len(encoding.encode(text))
 
def chunk_text(text, max_tokens=4000, model="gpt-4"):
   
    # Chunk text into smaller pieces to fit within the token limit.
    encoding = tiktoken.encoding_for_model(model)
    tokens = encoding.encode(text)
 
    # Split tokens into chunks of max_tokens size
    chunks = [tokens[i:i + max_tokens] for i in range(0, len(tokens), max_tokens)]
 
    # Convert token chunks back to text
    chunked_texts = [encoding.decode(chunk) for chunk in chunks]
    return chunked_texts
 
 
def get_s3_file_content(s3_url):
    try:
        # Determine the S3 bucket and object key
        if "s3.amazonaws.com" in s3_url:
            parts = s3_url.split('/')
            bucket_name = parts[2].split('.')[0]  # Extract bucket name
            object_key = '/'.join(parts[3:])     # Extract object key
        else:
            raise ValueError("Invalid S3 URL format")
        # Fetch the object from S3
        response = s3_client.get_object(Bucket=bucket_name, Key=object_key)
        content_type,content = response['ContentType'],response['Body']
       
        # Log for debugging purposes
        st.write(f"Bucket: {bucket_name}, Key: {object_key}, File Extension: {content_type}")
 
        if content_type == 'audio/mp3':
            local_path = r"C:\Users\visho\Documents\BDIAProjects\Assignment1_Copy\Assignment_1\A1\data\downloaded_file.mp3"
            s3_client.download_file(bucket_name, object_key,local_path)
            return extract_text_from_audio(local_path)
        
        elif content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return pd.read_excel(content.read())  # Handle .xlsx files
        elif content_type == 'application/pdf':
            return extract_text_from_pdf(content.read())
        elif content_type == 'application/zip':
            return extract_text_from_zip(content.read())
        elif content_type =='application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return extract_text_from_pptx(content.read())
        elif content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return extract_text_from_docx(content.read())
        elif content_type == 'image/png' or content_type == 'image/jpeg' or content_type == 'image/jpg':
            local_path = r"C:\Users\visho\Documents\BDIAProjects\Assignment1_Copy\Assignment_1\A1\data\downloaded_file" +"."+ str(content_type.split("/")[1])
            s3_client.download_file(bucket_name, object_key,local_path)
            return analyze_image(local_path)  # OCR for PNG/JPG
       
       
        else:
            # Try to decode as UTF-8 if it's a plain text file
            try:
                with open(local_path, "r", encoding="utf-8") as f:
                    return f.read()
            except UnicodeDecodeError:
                st.error(f"Cannot decode the file content at {s3_url} as text. It might be a binary file.")
                return None
               
 
    except Exception as e:
        st.error(f"Error processing S3 file: {e}")
        raise
 
def show():
    # Connect to Amazon RDS Database
    try:
        db = DBConnection.get_instance()
        cursor = db.get_cursor()
    except Exception as e:
        st.error(f"Error connecting to the database: {e}")
        return
 
    # Fetch data from the database
    cursor.execute('SELECT * FROM validation_table;')
    validation_table = cursor.fetchall()
    test_cases = pd.DataFrame(validation_table, columns=cursor.column_names)
 
    st.title("GAIA Benchmark Evaluation App")
 
    # Select the test case and store it in session state
    test_case_options = dict(zip(test_cases['serial_no'], test_cases['question']))
    selected_test_case = st.selectbox("Select a Test Case:", options=test_case_options.keys(), key="select_test_case")
 
    if selected_test_case:
        selected_question = test_case_options[selected_test_case]
        file_path_list = test_cases[test_cases['serial_no'] == selected_test_case]['file_path']  # Get file path
        file_path = None if file_path_list.empty else file_path_list.values[0]
       
        st.session_state['selected_test_case'] = selected_test_case  # Save selected test case in session state
        st.session_state['selected_question'] = selected_question    # Save selected question in session state
 
        # Display selected question
        st.write(f"Question: {selected_question}")
        if file_path:
            st.write(f"Accessing File: {file_path.split('/')[-1]}")  # Print the file name with extension
 
        # Search button - trigger OpenAI API call when clicked
            if st.button("Test",key='b1'):
                try:
                    # Access the file content from S3
                    file_content = get_s3_file_content(file_path)
                    st.session_state['file_content'] = file_content
 
                    if file_content is None:
                        # If file content could not be retrieved (e.g., binary file)
                        st.error("The file cannot be processed because it is not in a valid text format.")
                    elif isinstance(file_content, pd.DataFrame):
                        # If the file is an Excel file (DataFrame), limit the number of rows/columns to reduce size
                        st.write("Excel File Content (limited view):")
                        limited_content = file_content.head(10)  # Limit to the first 10 rows
                        st.dataframe(limited_content)  # Display the DataFrame in the Streamlit app
 
                        # Convert the limited DataFrame content to a string for the OpenAI API call
                        limited_content_str = limited_content.to_string()
 
                        # OpenAI API call - pass limited file content along with the question
                        response = client.chat.completions.create(
                            model="gpt-4",
                            messages=[
                                {"role": "user", "content": f"File content:\n{limited_content_str}\n\nQuestion: {selected_question}"}
                            ],
                            max_tokens=3000  # Adjust token limit for larger inputs
                        )
 
                        # Access the choices if available
                        if response.choices and len(response.choices) > 0:
                            openai_response_content = response.choices[0].message.content
                            st.session_state["openai_response"] = openai_response_content
                            st.success("Response from OpenAI successfully received!")
 
                            # Navigate to page 2 to display the response
                            st.session_state["page"] = "2_Test_Case"
                            st.rerun()  # This will trigger the app to rerun and go to page 2
                        else:
                            st.error("Unexpected API response format. No choices found.")
                    else:
                        # Handle plain text file content
                        token_count = count_tokens(file_content)
                        st.write(f"Token count: {token_count}")
 
                        if token_count > 7000:
                            st.warning("Content is too large. Chunking the content.")
                            chunks = chunk_text(file_content, max_tokens=4000)
                        else:
                            chunks = [file_content]
 
                        for chunk in chunks:
                            # OpenAI API call - pass chunked content along with the question
                            response = client.chat.completions.create(
                                model="gpt-4",
                                messages=[
                                    {"role": "user", "content": f"File content:\n{chunk}\n\nQuestion: {selected_question}"}
                                ],
                                max_tokens=3000  # Use a reasonable limit per request
                            )
 
                            if response.choices and len(response.choices) > 0:
                                openai_response_content = response.choices[0].message.content
                                st.session_state["openai_response"] = openai_response_content
                                st.success("Response from OpenAI successfully received!")
 
                                # Navigate to page 2 to display the response
                                st.session_state["page"] = "2_Test_Case"
                                st.rerun()
                            else:
                                st.error("Unexpected API response format. No choices found.")
 
                except Exception as e:
                    st.error(f"Error during OpenAI API call or S3 file access: {e}")
        else:
        # If no file path is available, answer based on the question only
            st.warning("No file path found. Answering based on the question only.")
            if st.button("Search"):
                try:
                    # OpenAI API call - pass the question directly without file content
                    response = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[
                            {"role": "user", "content": f"Question: {selected_question}"}
                        ],
                        max_tokens=3000  # Use a reasonable limit per request
                    )
 
                    # Safely access the choices if available
                    if response.choices and len(response.choices) > 0:
                        openai_response_content = response.choices[0].message.content
                        st.session_state["openai_response"] = openai_response_content
                        st.success("Response from OpenAI successfully received!")
 
                        # Navigate to page 2 to display the response
                        st.session_state["page"] = "2_Test_Case"
                        st.rerun()  # This will trigger the app to rerun and go to page 2
                    else:
                        st.error("Unexpected API response format. No choices found.")
                except Exception as e:
                    st.error(f"Error during OpenAI API call: {e}")
 
 
        # Add the "Go to Summary" button
        if st.button("Go to Summary"):
                st.session_state["page"] = "summary"
 